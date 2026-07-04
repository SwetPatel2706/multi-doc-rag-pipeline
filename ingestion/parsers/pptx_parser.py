import os
from typing import List, Dict, Any
from pptx import Presentation

def parse_pptx(pptx_path: str) -> List[Dict[str, Any]]:
    """
    Reads a PPTX slide deck and returns a list of dictionaries, one per slide.
    Format:
    {
        "text": str,
        "metadata": {
            "source": str (filename),
            "type": "pptx",
            "slide": int (1-based),
            "title": str
        }
    }
    """
    filename = os.path.basename(pptx_path)
    print(f"Parsing PPTX: {filename}")
    prs = Presentation(pptx_path)
    results = []

    for i, slide in enumerate(prs.slides, start=1):
        title = ""
        # 1. Try built-in title shape
        try:
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                title = slide.shapes.title.text.strip()
        except Exception:
            pass

        # 2. Fallback: scan placeholders for TITLE or CENTER_TITLE types
        if not title:
            for shape in slide.shapes:
                if shape.is_placeholder:
                    try:
                        ph_type = shape.placeholder_format.type
                        # 1: TITLE, 3: CENTER_TITLE
                        if ph_type in (1, 3) and shape.has_text_frame:
                            title = shape.text.strip()
                            if title:
                                break
                    except Exception:
                        pass

        # 3. Fallback: first non-empty short text shape
        if not title:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_val = shape.text.strip()
                    if text_val and len(text_val) < 100:
                        title = text_val
                        break

        # Collect all text on slide
        slide_text_parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_val = shape.text.strip()
                if text_val:
                    # Replace multiple whitespaces/newlines within paragraphs
                    cleaned_paragraphs = []
                    for para in shape.text_frame.paragraphs:
                        p_text = " ".join(para.text.split())
                        if p_text:
                            cleaned_paragraphs.append(p_text)
                    if cleaned_paragraphs:
                        slide_text_parts.append("\n".join(cleaned_paragraphs))

        slide_text = "\n\n".join(slide_text_parts)

        # Let's ensure the title is at least represented if found
        if not title:
            title = "Untitled Slide"

        results.append({
            "text": slide_text,
            "metadata": {
                "source": filename,
                "type": "pptx",
                "slide": i,
                "title": title
            }
        })

    print(f"  → Extracted {len(results)} slides from {filename}")
    return results
